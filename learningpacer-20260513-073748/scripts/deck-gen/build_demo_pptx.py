"""LearningPacer — PRODUCT DEMO deck (Keynote-grade rendering).
12 slides. Same Claude warm-paper palette, but every diagram has soft
shadows, refined proportions, native arcs / donuts / circles instead
of stacked ovals, and Apple-style breathing room.

Output: .local/fyp-deck/LearningPacer_Product_Demo_Deck.pptx
Run:    python3 scripts/deck-gen/build_demo_pptx.py
"""
import os, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── palette (Claude warm-paper, refined) ────────────────────────────────────
PAPER  = RGBColor(0xfa, 0xf6, 0xee)
PAPER2 = RGBColor(0xf3, 0xed, 0xdd)
INK    = RGBColor(0x1e, 0x29, 0x3b)
INK2   = RGBColor(0x47, 0x55, 0x69)
INK3   = RGBColor(0x64, 0x74, 0x8b)
MUTED  = RGBColor(0x94, 0xa3, 0xb8)
HAIR   = RGBColor(0xd4, 0xd9, 0xe0)   # softer, less heavy than slate-300
HAIR2  = RGBColor(0xe5, 0xe7, 0xeb)
TEAL   = RGBColor(0x0f, 0x76, 0x6e)
TEAL_D = RGBColor(0x0a, 0x52, 0x4d)
TEAL_L = RGBColor(0xd4, 0xe7, 0xe5)   # very light teal tint
CYAN   = RGBColor(0x08, 0x91, 0xb2)
AMBER  = RGBColor(0xb4, 0x53, 0x09)
RED    = RGBColor(0xb9, 0x1c, 0x1c)
SOFT   = RGBColor(0xef, 0xe9, 0xda)
SOFT2  = RGBColor(0xe6, 0xdf, 0xcc)
SOFT3  = RGBColor(0xf6, 0xf1, 0xe2)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GREY   = RGBColor(0xeb, 0xed, 0xf0)

SERIF_BOLD = "Cambria"
SANS       = "Calibri"
SANS_BOLD  = "Calibri"
MONO       = "Consolas"

# ── canvas ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

OUT_DIR  = ".local/fyp-deck"
os.makedirs(OUT_DIR, exist_ok=True)
OUT_PATH = os.path.join(OUT_DIR, "LearningPacer_Product_Demo_Deck.pptx")
BLANK = prs.slide_layouts[6]

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ── primitives ──────────────────────────────────────────────────────────────
def _no_outline(shape): shape.line.fill.background()
def fill(shape, color): shape.fill.solid(); shape.fill.fore_color.rgb = color
def stroke(shape, color, w=0.6):
    shape.line.color.rgb = color; shape.line.width = Pt(w)


def add_shadow(shape, *, blur=10, dist=3, alpha=14, angle_deg=90):
    """Apply a soft outer drop-shadow via raw OOXML.
    blur/dist in points; alpha 0-100 (pct); angle 90 = down."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    blur_emu = int(blur * 12700)
    dist_emu = int(dist * 12700)
    angle_u  = int(angle_deg * 60000)
    alpha_u  = int(alpha * 1000)
    xml = (
        f'<a:effectLst xmlns:a="{NSMAP_A}">'
        f'<a:outerShdw blurRad="{blur_emu}" dist="{dist_emu}" '
        f'dir="{angle_u}" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="{alpha_u}"/></a:srgbClr>'
        f'</a:outerShdw>'
        f'</a:effectLst>'
    )
    spPr.append(etree.fromstring(xml))


def rect(slide, x, y, w, h, *, fill_c=None, stroke_c=None, stroke_w=0.6,
         rounded=False, radius=0.08):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, x, y, w, h)
    s.shadow.inherit = False
    if fill_c is not None: fill(s, fill_c)
    else:                  s.fill.background()
    if stroke_c is not None: stroke(s, stroke_c, stroke_w)
    else:                    _no_outline(s)
    if rounded:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s


def oval(slide, cx, cy, r, *, fill_c=None, stroke_c=None, stroke_w=0.6):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    s.shadow.inherit = False
    if fill_c is not None: fill(s, fill_c)
    else:                  s.fill.background()
    if stroke_c is not None: stroke(s, stroke_c, stroke_w)
    else:                    _no_outline(s)
    return s


def donut(slide, cx, cy, r, *, fill_c, hole_frac=0.62):
    """Native donut — ring shape, used for the progress dial."""
    s = slide.shapes.add_shape(MSO_SHAPE.DONUT,
                                cx - r, cy - r, r * 2, r * 2)
    s.shadow.inherit = False
    fill(s, fill_c); _no_outline(s)
    try:
        # adjustment is the inner-radius fraction (0..0.5)
        s.adjustments[0] = max(0.05, min(0.49, (1 - hole_frac) / 2))
    except Exception:
        pass
    return s


def line(slide, x1, y1, x2, y2, *, color=HAIR, w=0.4, dash=False):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    s.line.color.rgb = color
    s.line.width = Pt(w)
    if dash:
        ln = s.line._get_or_add_ln()
        d = etree.SubElement(ln, qn("a:prstDash")); d.set("val", "dash")
    return s


def arrow(slide, x1, y1, x2, y2, *, color=TEAL, w=0.9, dash=False,
          head="triangle"):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    s.line.color.rgb = color
    s.line.width = Pt(w)
    ln = s.line._get_or_add_ln()
    if dash:
        d = etree.SubElement(ln, qn("a:prstDash")); d.set("val", "dash")
    t = etree.SubElement(ln, qn("a:tailEnd"))
    t.set("type", head); t.set("w", "med"); t.set("len", "med")
    return s


def textbox(slide, x, y, w, h, text, *, font=SANS, size=12, bold=False,
            italic=False, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
            char_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color
        if char_spacing is not None:
            # subtle letter-spacing via raw XML on rPr
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def rich(slide, x, y, w, h, paragraphs, *, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, st in runs:
            r = p.add_run(); r.text = txt
            r.font.name  = st.get("font", SANS)
            r.font.size  = Pt(st.get("size", 12))
            r.font.bold  = st.get("bold", False)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", INK)
    return tb


def paper(slide, color=PAPER):
    bg = rect(slide, 0, 0, SW, SH, fill_c=color); _no_outline(bg)
    return bg


def card(slide, x, y, w, h, *, fill_c=PAPER, stroke_c=HAIR, stroke_w=0.45,
         radius=0.04, shadow=True, blur=14, dist=4, alpha=10):
    """Standard 'Keynote' card — soft outline + subtle drop shadow."""
    c = rect(slide, x, y, w, h, fill_c=fill_c, stroke_c=stroke_c,
             stroke_w=stroke_w, rounded=True, radius=radius)
    if shadow:
        add_shadow(c, blur=blur, dist=dist, alpha=alpha)
    return c


# ── chrome ──────────────────────────────────────────────────────────────────
def chrome(slide, page, total, label):
    paper(slide)
    # tiny teal corner mark
    rect(slide, Inches(0.55), Inches(0.45), Inches(0.22), Inches(0.22),
         fill_c=TEAL)
    rect(slide, Inches(0.55), Inches(0.55), Inches(0.22), Inches(0.12),
         fill_c=TEAL_D)
    textbox(slide, Inches(0.95), Inches(0.45), Inches(8), Inches(0.3),
            label.upper(),
            font=SANS_BOLD, size=10, bold=True, color=TEAL,
            line_spacing=1.0, char_spacing=120)
    textbox(slide, Inches(8.5), Inches(0.45), Inches(4.3), Inches(0.3),
            "LEARNINGPACER  ·  PRODUCT DEMO",
            font=SANS_BOLD, size=9, bold=True, color=MUTED,
            align=PP_ALIGN.RIGHT, char_spacing=120)
    line(slide, Inches(0.55), Inches(0.92), Inches(12.78), Inches(0.92),
         color=HAIR, w=0.4)
    line(slide, Inches(0.55), Inches(7.05), Inches(12.78), Inches(7.05),
         color=HAIR, w=0.4)
    textbox(slide, Inches(0.55), Inches(7.15), Inches(8), Inches(0.3),
            "learningpacer.app  ·  ELEC3120 Computer Networks",
            font=MONO, size=8, color=MUTED)
    textbox(slide, Inches(8.5), Inches(7.15), Inches(4.3), Inches(0.3),
            f"{page:02d} / {total:02d}",
            font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def big_title(slide, eyebrow, title, sub=None, y=Inches(1.2)):
    textbox(slide, Inches(0.55), y, Inches(12), Inches(0.32),
            eyebrow.upper(),
            font=SANS_BOLD, size=10, bold=True, color=TEAL,
            char_spacing=180)
    textbox(slide, Inches(0.55), y + Inches(0.32), Inches(12), Inches(0.95),
            title, font=SERIF_BOLD, size=36, bold=True, color=INK,
            line_spacing=1.05)
    if sub:
        textbox(slide, Inches(0.55), y + Inches(1.35), Inches(12), Inches(0.5),
                sub, font="Cambria", size=15, italic=True, color=INK2,
                line_spacing=1.3)


# ─────────────────────────────────────────────────────────────────────────────
# REFINED DRAWING HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def signal_mark(slide, x, y, w, h):
    """5 vertical bars rising teal. Rounded ends + subtle shadow + accent bar."""
    n = 5
    gap = w * 0.06
    bw = (w - gap * (n - 1)) / n
    for i in range(n):
        bh = h * (0.22 + 0.165 * i)
        bx = x + i * (bw + gap)
        by = y + (h - bh)
        s = rect(slide, bx, by, bw, bh, fill_c=TEAL,
                 rounded=True, radius=0.5)
        add_shadow(s, blur=6, dist=1.5, alpha=18)
    # underline tick
    line(slide, x, y + h + Inches(0.08), x + w, y + h + Inches(0.08),
         color=TEAL_D, w=2.4)


def browser_chrome(slide, x, y, w, h, url):
    """Drawn browser window — body + chrome bar + traffic lights + URL pill.
    Ships with a soft drop shadow, like a Keynote app mockup."""
    body = rect(slide, x, y, w, h, fill_c=WHITE, stroke_c=HAIR, stroke_w=0.5,
                rounded=True, radius=0.035)
    add_shadow(body, blur=22, dist=8, alpha=18)
    bar = rect(slide, x, y, w, Inches(0.42), fill_c=SOFT3,
               rounded=True, radius=0.045)
    rect(slide, x, y + Inches(0.20), w, Inches(0.22), fill_c=SOFT3)
    # traffic lights (with subtle inner highlight ring)
    lights = [
        (RGBColor(0xff, 0x5f, 0x57), RGBColor(0xe0, 0x44, 0x3e)),
        (RGBColor(0xfe, 0xbc, 0x2e), RGBColor(0xdf, 0x9f, 0x12)),
        (RGBColor(0x28, 0xc8, 0x40), RGBColor(0x1a, 0xa6, 0x2f)),
    ]
    for i, (top_c, edge_c) in enumerate(lights):
        cx_ = x + Inches(0.22 + i * 0.22)
        cy_ = y + Inches(0.21)
        oval(slide, cx_, cy_, Inches(0.075), fill_c=edge_c)
        oval(slide, cx_, cy_, Inches(0.063), fill_c=top_c)
    # URL pill
    pill_w = w - Inches(1.7)
    pill_x = x + Inches(1.15)
    pill = rect(slide, pill_x, y + Inches(0.105), pill_w, Inches(0.21),
                fill_c=PAPER, stroke_c=HAIR, stroke_w=0.4,
                rounded=True, radius=0.5)
    # tiny lock glyph
    lock_x = pill_x + Inches(0.14)
    lock_y = y + Inches(0.155)
    rect(slide, lock_x, lock_y + Inches(0.045), Inches(0.07),
         Inches(0.06), fill_c=INK2, rounded=True, radius=0.18)
    # url text
    textbox(slide, pill_x + Inches(0.30), y + Inches(0.105),
            pill_w - Inches(0.4), Inches(0.21), url,
            font=MONO, size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    line(slide, x, y + Inches(0.42), x + w, y + Inches(0.42),
         color=HAIR2, w=0.4)
    return (x, y + Inches(0.42), w, h - Inches(0.42))


def chat_bubble(slide, x, y, w, *, text, from_user=False, h=None,
                cite=None, font_size=11, shadow=True):
    body_h = h or Inches(0.7)
    if from_user:
        bg = rect(slide, x, y, w, body_h, fill_c=TEAL,
                  rounded=True, radius=0.20)
        if shadow:
            add_shadow(bg, blur=10, dist=2.5, alpha=22)
        textbox(slide, x + Inches(0.22), y + Inches(0.14),
                w - Inches(0.44), body_h - Inches(0.28), text,
                font=SANS, size=font_size, color=WHITE, line_spacing=1.4)
    else:
        bg = rect(slide, x, y, w, body_h, fill_c=PAPER,
                  stroke_c=HAIR, stroke_w=0.45,
                  rounded=True, radius=0.13)
        if shadow:
            add_shadow(bg, blur=14, dist=3, alpha=12)
        # subtle left accent
        rect(slide, x, y + Inches(0.18), Inches(0.05),
             body_h - Inches(0.36), fill_c=TEAL,
             rounded=True, radius=0.5)
        textbox(slide, x + Inches(0.24), y + Inches(0.14),
                w - Inches(0.48), body_h - Inches(0.28), text,
                font=SANS, size=font_size, color=INK, line_spacing=1.4)
        if cite:
            cy_ = y + body_h - Inches(0.36)
            cb = rect(slide, x + Inches(0.24), cy_, Inches(2.2),
                      Inches(0.24), fill_c=TEAL_L,
                      stroke_c=TEAL, stroke_w=0.55,
                      rounded=True, radius=0.5)
            textbox(slide, x + Inches(0.24), cy_, Inches(2.2),
                    Inches(0.24), cite,
                    font=MONO, size=8, bold=True, color=TEAL,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return (x, y, w, body_h)


def callout(slide, anchor_x, anchor_y, label_x, label_y, label_w, label_h,
            label, *, color=AMBER):
    """Annotation callout: dot + dashed leader + label box."""
    oval(slide, anchor_x, anchor_y, Inches(0.075), fill_c=color)
    oval(slide, anchor_x, anchor_y, Inches(0.038), fill_c=PAPER)
    line(slide, anchor_x, anchor_y, label_x, label_y + label_h / 2,
         color=color, w=0.7, dash=True)
    box = rect(slide, label_x, label_y, label_w, label_h,
               fill_c=PAPER, stroke_c=color, stroke_w=0.7,
               rounded=True, radius=0.18)
    add_shadow(box, blur=10, dist=2, alpha=8)
    textbox(slide, label_x + Inches(0.14), label_y + Inches(0.09),
            label_w - Inches(0.24), label_h - Inches(0.18), label,
            font=SANS, size=10, color=INK, line_spacing=1.32)


def goal_ring(slide, cx, cy, r, percent, *,
              ring_thickness=0.18, color=TEAL, track=SOFT2):
    """Apple-Activity-style goal dial: light track ring + a single
    position marker dot rotated to the `percent` angle (clockwise from
    12 o'clock). Center is paper so the % label sits cleanly.

    Note: PPTX has no robust partial-arc primitive that survives across
    PowerPoint / Keynote / Google Slides — a marker dot reads more
    reliably than a stitched-arc kludge."""
    # background track
    donut(slide, cx, cy, r, fill_c=track, hole_frac=1 - 2 * ring_thickness)
    # inner paper face
    inner_r = int(r * (1 - 2 * ring_thickness)) - Inches(0.02)
    oval(slide, cx, cy, inner_r, fill_c=PAPER)
    # tiny start dot at 12 o'clock
    start_dot = oval(slide, cx, cy - int(r * (1 - ring_thickness)),
                     Inches(0.045), fill_c=color)
    # marker dot at percentage angle (clockwise from top)
    theta = math.radians(-90 + (percent / 100.0) * 360.0)
    mid_r = int(r * (1 - ring_thickness))
    mx = int(cx + math.cos(theta) * mid_r)
    my = int(cy + math.sin(theta) * mid_r)
    marker_r = int(r * ring_thickness * 0.95)
    marker = oval(slide, mx, my, marker_r,
                  fill_c=color, stroke_c=PAPER, stroke_w=1.0)
    add_shadow(marker, blur=6, dist=1.5, alpha=24)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    paper(s)

    # left column
    signal_mark(s, Inches(0.7), Inches(0.85), Inches(1.05), Inches(0.75))
    textbox(s, Inches(0.7), Inches(1.92), Inches(8), Inches(0.4),
            "LEARNINGPACER",
            font=SANS_BOLD, size=14, bold=True, color=TEAL,
            char_spacing=200)

    textbox(s, Inches(0.7), Inches(2.55), Inches(8.5), Inches(2.7),
            "A virtual TA\nthat actually\nknows the slides.",
            font=SERIF_BOLD, size=64, bold=True, color=INK,
            line_spacing=1.02)

    line(s, Inches(0.7), Inches(5.55), Inches(5.5), Inches(5.55),
         color=TEAL, w=2.0)

    textbox(s, Inches(0.7), Inches(5.78), Inches(8.5), Inches(0.55),
            "Eight-minute product walk-through.",
            font="Cambria", size=18, italic=True, color=INK2)

    rich(s, Inches(0.7), Inches(6.65), Inches(8), Inches(0.4),
         [[("DEMO   ",
            {"bold": True, "color": TEAL, "size": 10}),
           ("HKUST · ELEC3120 Computer Networks · May 2026",
            {"color": INK2, "size": 11})]])

    # right column: full-bleed teal panel
    panel = rect(s, Inches(8.95), 0, SW - Inches(8.95), SH, fill_c=TEAL)
    _no_outline(panel)
    # inner darker layer
    inner = rect(s, Inches(9.25), Inches(0.6), Inches(3.85), Inches(6.3),
                 fill_c=TEAL_D, rounded=True, radius=0.04)
    add_shadow(inner, blur=20, dist=0, alpha=20)
    # subtle hairline highlight at inner card top edge
    line(s, Inches(9.35), Inches(0.65), Inches(13.05), Inches(0.65),
         color=RGBColor(0x35, 0x91, 0x88), w=0.6)

    textbox(s, Inches(9.45), Inches(0.85), Inches(3.5), Inches(0.4),
            "WHAT YOU'LL SEE",
            font=SANS_BOLD, size=10, bold=True, color=PAPER,
            char_spacing=200)

    pillars = [
        ("01", "A real answer",        "Cited from a real lecture slide."),
        ("02", "A real mock paper",    "Generated live, in front of you."),
        ("03", "A real study session", "Quiz · Pomodoro · weak-topic stats."),
    ]
    py = Inches(1.4); ph = Inches(1.65)
    for i, (n, t, d) in enumerate(pillars):
        cy = py + i * ph
        textbox(s, Inches(9.45), cy, Inches(3.5), Inches(0.5), n,
                font=SERIF_BOLD, size=24, bold=True, color=PAPER,
                line_spacing=1.0)
        line(s, Inches(9.45), cy + Inches(0.55),
             Inches(10.4), cy + Inches(0.55), color=PAPER, w=1.5)
        textbox(s, Inches(9.45), cy + Inches(0.65), Inches(3.5),
                Inches(0.45), t,
                font=SANS_BOLD, size=14, bold=True, color=PAPER)
        textbox(s, Inches(9.45), cy + Inches(1.10), Inches(3.5),
                Inches(0.55), d,
                font=SANS, size=10.5, color=PAPER, line_spacing=1.35)

    textbox(s, Inches(9.45), Inches(6.6), Inches(3.5), Inches(0.4),
            "12 SLIDES  ·  8 MINUTES",
            font=MONO, size=9, color=PAPER, char_spacing=200)
    return s


def slide_moment(page, total):
    s = prs.slides.add_slide(BLANK)
    paper(s)
    textbox(s, Inches(0.55), Inches(0.55), Inches(8), Inches(0.3),
            "01  ·  THE MOMENT",
            font=SANS_BOLD, size=10, bold=True, color=TEAL,
            char_spacing=200)
    textbox(s, Inches(8.5), Inches(0.55), Inches(4.3), Inches(0.3),
            f"{page:02d} / {total:02d}",
            font=SANS, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    textbox(s, Inches(0.7), Inches(1.95), Inches(2), Inches(0.55),
            "\u201c", font=SERIF_BOLD, size=160, bold=True, color=TEAL_L,
            line_spacing=0.8)
    textbox(s, Inches(1.7), Inches(2.4), Inches(11), Inches(3.6),
            "It's 2 a.m. The exam\nis in nine hours.\nThere's no one to ask.",
            font=SERIF_BOLD, size=58, bold=True, color=INK,
            line_spacing=1.05)
    line(s, Inches(1.7), Inches(6.2), Inches(5.5), Inches(6.2),
         color=TEAL, w=1.8)
    textbox(s, Inches(1.7), Inches(6.35), Inches(11), Inches(0.4),
            "— What every ELEC3120 student is doing the night before the final.",
            font="Cambria", size=14, italic=True, color=INK2)
    return s


def slide_meet(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "02  ·  Meet LearningPacer")
    big_title(s, "the product",
              "One companion, three jobs.",
              "Tutor when you're stuck. Examiner when you're ready. Coach the rest of the time.")

    cards = [
        ("Tutor",
         "Conversation that cites the slide.",
         ["Lecture-grounded answers",
          "EN / 中文 inline",
          "Refuses out-of-scope cleanly"]),
        ("Examiner",
         "Generates lecturer-grade papers.",
         ["Mock midterm + final on demand",
          "Trap-analysis booklet",
          "Print-ready PDFs"]),
        ("Coach",
         "Quiz · Pomodoro · weak-topic stats.",
         ["Tracks every minute studied",
          "Daily goal ring",
          "Spaced-repetition flashcards"]),
    ]
    base_x = Inches(0.55); base_y = Inches(3.4)
    cw = Inches(4.05); ch = Inches(3.25); gap = Inches(0.07)
    for i, (t, sub, items) in enumerate(cards):
        cx = base_x + i * (cw + gap)
        c = card(s, cx, base_y, cw, ch, fill_c=PAPER, blur=18, dist=5)
        # left teal accent column
        rect(s, cx, base_y, Inches(0.18), ch, fill_c=TEAL,
             rounded=True, radius=0.5)
        rect(s, cx + Inches(0.05), base_y, Inches(0.13), ch, fill_c=TEAL)
        textbox(s, cx + Inches(0.45), base_y + Inches(0.32),
                cw - Inches(0.6), Inches(0.4), t,
                font=SERIF_BOLD, size=24, bold=True, color=INK)
        textbox(s, cx + Inches(0.45), base_y + Inches(0.92),
                cw - Inches(0.6), Inches(0.4), sub,
                font="Cambria", size=12, italic=True, color=INK2)
        line(s, cx + Inches(0.45), base_y + Inches(1.4),
             cx + Inches(1.4), base_y + Inches(1.4),
             color=TEAL, w=1.4)
        for j, it in enumerate(items):
            iy = base_y + Inches(1.62 + j * 0.46)
            oval(s, cx + Inches(0.55), iy + Inches(0.13),
                 Inches(0.045), fill_c=TEAL)
            textbox(s, cx + Inches(0.75), iy, cw - Inches(0.9),
                    Inches(0.42), it,
                    font=SANS, size=11, color=INK)
    return s


def slide_anatomy(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "03  ·  Anatomy of an answer")
    big_title(s, "how it differs",
              "Every answer carries its source.")

    bx, by, bw, bh = Inches(2.65), Inches(3.0), Inches(8.05), Inches(3.95)
    inner = browser_chrome(s, bx, by, bw, bh, "learningpacer.app/tutor")
    cx = inner[0] + Inches(0.30)
    cy = inner[1] + Inches(0.22)
    cw_ = inner[2] - Inches(0.6)

    user_b = chat_bubble(s, cx + Inches(2.2), cy, cw_ - Inches(2.2),
                         text=("Why does TCP Reno cut cwnd in half on a "
                               "triple-duplicate ACK? Cite the slide."),
                         from_user=True, h=Inches(0.7))
    a_y = cy + Inches(0.95)
    asst_h = Inches(2.45)
    chat_bubble(s, cx, a_y, cw_,
                text=("Reno treats a triple-duplicate ACK as mild congestion: "
                      "the receiver is still ACK-ing, so the pipe is not empty. "
                      "It halves cwnd (multiplicative decrease) and re-enters "
                      "congestion-avoidance — not slow-start."),
                cite="Lec 09 · slide 17",
                h=asst_h)

    callout(s,
            anchor_x=cx + Inches(0.55), anchor_y=a_y + asst_h - Inches(0.24),
            label_x=Inches(0.55),       label_y=Inches(5.55),
            label_w=Inches(2.0),        label_h=Inches(0.88),
            label="Citation pill links to the exact lecture slide.",
            color=TEAL)
    callout(s,
            anchor_x=cx + cw_ - Inches(0.4), anchor_y=a_y + Inches(0.30),
            label_x=Inches(10.85),       label_y=Inches(3.0),
            label_w=Inches(2.0),         label_h=Inches(0.88),
            label="Lecturer's vocabulary, not generic networking prose.",
            color=AMBER)
    callout(s,
            anchor_x=cx + Inches(2.2), anchor_y=a_y + Inches(1.6),
            label_x=Inches(10.85),     label_y=Inches(4.1),
            label_w=Inches(2.0),       label_h=Inches(0.88),
            label="Distinguishes Reno's MD from Tahoe's slow-start reset.",
            color=CYAN)
    callout(s,
            anchor_x=user_b[0] + Inches(0.6),
            anchor_y=user_b[1] + Inches(0.35),
            label_x=Inches(10.85), label_y=Inches(5.2),
            label_w=Inches(2.0),    label_h=Inches(0.88),
            label="User can demand a citation — the model has to comply.",
            color=TEAL_D)

    textbox(s, Inches(0.55), Inches(6.65), Inches(12), Inches(0.3),
            "Same answer in EN / 中文; refuses cleanly when no slide matches.",
            font="Cambria", size=11, italic=True, color=INK2)
    return s


def slide_pipeline(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "04  ·  Under the hood")
    big_title(s, "the pipeline",
              "Question to cited answer in five steps.",
              "Same flow every time — no model freelancing.")

    stages = [
        ("ASK",      "Student types\nin EN or 中文"),
        ("EMBED",    "Question →\nvector"),
        ("RETRIEVE", "Top-k slide\nchunks"),
        ("REASON",   "Claude\n+ tool belt"),
        ("CITE",     "Answer + slide pill"),
    ]
    n = len(stages)
    base_x = Inches(0.85); right_x = Inches(12.48)
    base_y = Inches(4.05)
    radius = Inches(0.55)

    cx_pos = [base_x + (right_x - base_x) / (n - 1) * i for i in range(n)]

    # rail under the circles, behind everything
    line(s, cx_pos[0], base_y, cx_pos[-1], base_y,
         color=HAIR2, w=2.4)

    for i, (cap, body) in enumerate(stages):
        cx_ = cx_pos[i]
        # halo (ultra-soft outer ring)
        oval(s, cx_, base_y, radius + Inches(0.13),
             fill_c=TEAL_L)
        # main circle
        circ = oval(s, cx_, base_y, radius, fill_c=PAPER,
                    stroke_c=TEAL, stroke_w=1.6)
        add_shadow(circ, blur=14, dist=4, alpha=18)
        # number inside
        textbox(s, cx_ - radius, base_y - Inches(0.32),
                radius * 2, Inches(0.45), f"{i+1:02d}",
                font=SERIF_BOLD, size=22, bold=True, color=TEAL,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, cx_ - radius, base_y + Inches(0.05),
                radius * 2, Inches(0.3), cap,
                font=SANS_BOLD, size=8.5, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                char_spacing=120)
        # body caption below the circle
        textbox(s, cx_ - Inches(1.05), base_y + radius + Inches(0.32),
                Inches(2.1), Inches(0.95), body,
                font=SANS, size=10.5, color=INK2,
                align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Trust band
    band_y = Inches(5.95)
    band = card(s, Inches(0.55), band_y, Inches(12.23), Inches(0.85),
                fill_c=SOFT3, stroke_c=HAIR, stroke_w=0.4,
                blur=12, dist=3, alpha=8, radius=0.05)
    rich(s, Inches(0.85), band_y + Inches(0.18),
         Inches(11.63), Inches(0.55),
         [[("Why this matters.  ",
            {"bold": True, "color": TEAL, "size": 12}),
           ("If retrieval returns nothing relevant, the model is told to "
            "refuse — not to invent. Hallucinations don't make it past "
            "the retrieval check.",
            {"color": INK, "size": 12})]],
         line_spacing=1.4)

    textbox(s, Inches(0.55), Inches(6.85), Inches(12.23), Inches(0.18),
            "ANTHROPIC CLAUDE SONNET 4   ·   OPENROUTER FALLBACK FLEET   ·   "
            "SUPABASE PGVECTOR   ·   NEXT.JS 16 EDGE ROUTES",
            font=MONO, size=8, color=MUTED, align=PP_ALIGN.CENTER,
            char_spacing=180)
    return s


def slide_moments(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "05  ·  Five product moments")
    big_title(s, "the surface",
              "What you'll click in the next 8 minutes.")

    base_x = Inches(0.55); base_y = Inches(3.3)
    span_w = Inches(12.23)
    n = 5; gap = Inches(0.18)
    cw = (span_w - (n - 1) * gap) / n
    ch = Inches(3.55)

    tiles = [
        ("Tutor",       "Ask. Get an answer. Click the slide cite.", "tutor"),
        ("Mock paper",  "One prompt → 33-page midterm + final.",     "exam"),
        ("Quiz",        "5-question lecture quiz with explanations.", "quiz"),
        ("Stats ring",  "Daily-minute goal rendered as a ring.",      "ring"),
        ("Subnet calc", "Worked steps, not just the answer.",         "subnet"),
    ]
    for i, (title, sub, kind) in enumerate(tiles):
        x = base_x + i * (cw + gap)
        # outer card with soft shadow
        c = card(s, x, base_y, cw, ch, fill_c=PAPER,
                 blur=18, dist=5, alpha=14)
        # mini UI canvas — top half (subtle inner panel)
        ux = x + Inches(0.18)
        uy = base_y + Inches(0.22)
        uw = cw - Inches(0.36)
        uh = Inches(2.05)
        inner = rect(s, ux, uy, uw, uh, fill_c=SOFT3,
                     stroke_c=HAIR, stroke_w=0.35,
                     rounded=True, radius=0.05)
        _draw_mini_ui(s, ux, uy, uw, uh, kind)
        # caption
        textbox(s, x + Inches(0.18), base_y + Inches(2.45),
                cw - Inches(0.36), Inches(0.4), title,
                font=SANS_BOLD, size=14, bold=True, color=TEAL)
        textbox(s, x + Inches(0.18), base_y + Inches(2.85),
                cw - Inches(0.36), Inches(0.65), sub,
                font=SANS, size=10.5, color=INK2, line_spacing=1.35)
    return s


def _draw_mini_ui(slide, x, y, w, h, kind):
    """Polished native sketches for the five tiles."""
    pad = Inches(0.18)
    if kind == "tutor":
        # user bubble (small, top-right)
        chat_bubble(slide, x + Inches(0.7), y + Inches(0.18),
                    w - Inches(0.85),
                    text="Why does Reno halve cwnd?",
                    from_user=True, h=Inches(0.46),
                    font_size=8, shadow=False)
        # asst bubble (larger, bottom)
        chat_bubble(slide, x + Inches(0.18), y + Inches(0.8),
                    w - Inches(0.55),
                    text="Triple-dup ACK = mild congestion → "
                         "halve cwnd (factor ½), re-enter cong-avoid.",
                    cite="Lec 09 · 17",
                    h=Inches(1.05),
                    font_size=8, shadow=False)
    elif kind == "exam":
        # Stack of pages (3 layered rectangles with light shadow)
        for i, off in enumerate([(0.32, 0.30), (0.20, 0.18), (0.08, 0.08)]):
            ox, oy = off
            page_rect = rect(slide,
                             x + Inches(ox), y + Inches(oy),
                             w - Inches(0.85), h - Inches(0.55),
                             fill_c=WHITE, stroke_c=HAIR, stroke_w=0.35,
                             rounded=True, radius=0.04)
            if i == 2:
                add_shadow(page_rect, blur=10, dist=2.5, alpha=14)
        # Header on top page
        top_x = x + Inches(0.08); top_y = y + Inches(0.08)
        top_w = w - Inches(0.85)
        textbox(slide, top_x + Inches(0.18), top_y + Inches(0.12),
                top_w - Inches(0.36), Inches(0.22),
                "ELEC3120  ·  Mock Final",
                font=MONO, size=7.5, color=TEAL, char_spacing=120)
        line(slide, top_x + Inches(0.18), top_y + Inches(0.36),
             top_x + top_w - Inches(0.18), top_y + Inches(0.36),
             color=HAIR2, w=0.5)
        # Question and answer lines
        textbox(slide, top_x + Inches(0.18), top_y + Inches(0.42),
                top_w - Inches(0.36), Inches(0.22),
                "Q1.  Subnet 192.168.1.0 / 26",
                font=SANS_BOLD, size=8, bold=True, color=INK)
        for i in range(5):
            ly = top_y + Inches(0.7 + i * 0.18)
            line(slide, top_x + Inches(0.20), ly,
                 top_x + top_w - Inches(0.20), ly,
                 color=HAIR2, w=0.4)
    elif kind == "quiz":
        # MCQ
        textbox(slide, x + pad, y + pad,
                w - pad * 2, Inches(0.34),
                "OSI layer for IP?",
                font=SANS_BOLD, size=10.5, bold=True, color=INK)
        line(slide, x + pad, y + pad + Inches(0.4),
             x + w - pad, y + pad + Inches(0.4),
             color=HAIR2, w=0.4)
        labels = [("A", "Transport"),
                  ("B", "Network"),
                  ("C", "Data link"),
                  ("D", "Session")]
        for i, (lbl, txt) in enumerate(labels):
            row_y = y + pad + Inches(0.5 + i * 0.34)
            selected = (i == 1)
            # radio
            oval(slide, x + pad + Inches(0.13),
                 row_y + Inches(0.13), Inches(0.085),
                 fill_c=TEAL if selected else PAPER,
                 stroke_c=TEAL, stroke_w=0.7)
            if selected:
                oval(slide, x + pad + Inches(0.13),
                     row_y + Inches(0.13), Inches(0.035),
                     fill_c=PAPER)
            # label letter + text
            textbox(slide, x + pad + Inches(0.32), row_y,
                    Inches(0.18), Inches(0.28), lbl + ".",
                    font=SANS_BOLD, size=8.5, bold=True,
                    color=TEAL if selected else INK2)
            textbox(slide, x + pad + Inches(0.55), row_y,
                    w - pad * 2 - Inches(0.55), Inches(0.28), txt,
                    font=SANS, size=9,
                    color=INK if selected else INK2)
    elif kind == "ring":
        cx_ = x + w / 2
        cy_ = y + h / 2 + Inches(0.05)
        r   = Inches(0.78)
        goal_ring(slide, cx_, cy_, r, percent=72,
                  ring_thickness=0.16)
        textbox(slide, cx_ - r, cy_ - Inches(0.28),
                r * 2, Inches(0.5), "72%",
                font=SERIF_BOLD, size=26, bold=True, color=INK,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, cx_ - r, cy_ + Inches(0.10),
                r * 2, Inches(0.22), "of daily goal",
                font=SANS, size=8, color=MUTED, align=PP_ALIGN.CENTER)
    elif kind == "subnet":
        # IP + slash + result
        textbox(slide, x + pad, y + pad,
                w - pad * 2, Inches(0.32),
                "192.168.1.0 / 26",
                font=MONO, size=13, bold=True, color=INK)
        line(slide, x + pad, y + pad + Inches(0.46),
             x + w - pad, y + pad + Inches(0.46),
             color=HAIR2, w=0.4)
        rows = [("hosts",     "62"),
                ("net mask",  "255.255.255.192"),
                ("broadcast", ".63")]
        for i, (k, v) in enumerate(rows):
            ry = y + pad + Inches(0.6 + i * 0.36)
            textbox(slide, x + pad, ry, w / 2 - pad,
                    Inches(0.28), k,
                    font=SANS, size=9, color=INK2)
            textbox(slide, x + w / 2, ry, w / 2 - pad,
                    Inches(0.28), v,
                    font=MONO, size=9.5, bold=True, color=INK,
                    align=PP_ALIGN.RIGHT)


def slide_numbers(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "06  ·  Numbers that matter")
    big_title(s, "we measured",
              "We did not ship on vibes.",
              "Quantitative checks panellists can re-run on the spot.")

    base_x = Inches(0.55); base_y = Inches(3.5)
    span_w = Inches(12.23); ch = Inches(3.0)
    n = 4; gap = Inches(0.15)
    cw = (span_w - (n - 1) * gap) / n
    items = [
        ("30 / 30", "Protocol facts correct",
         "Past-paper questions, citation-grounded answers — vs ChatGPT 24/30."),
        ("21 / 25", "Inter-AI agreement",
         "On the 4 disagreements, our answer matches the lecture slide."),
        ("EN / 中文", "Bilingual inline",
         "Switches mid-sentence; matches how Cantonese students think."),
        ("0", "Hallucinated citations",
         "If retrieval returns nothing, the model is told to refuse."),
    ]
    for i, (n_, lbl, desc) in enumerate(items):
        x = base_x + i * (cw + gap)
        c = card(s, x, base_y, cw, ch, fill_c=PAPER,
                 blur=18, dist=5, alpha=12)
        # top teal hairline accent
        rect(s, x, base_y, cw, Inches(0.06), fill_c=TEAL,
             rounded=True, radius=0.5)
        # auto-shrink long stat strings
        sz = 60 if len(n_) <= 5 else 38
        textbox(s, x + Inches(0.3), base_y + Inches(0.4),
                cw - Inches(0.6), Inches(1.1), n_,
                font=SERIF_BOLD, size=sz, bold=True, color=TEAL,
                line_spacing=1.0)
        textbox(s, x + Inches(0.3), base_y + Inches(1.6),
                cw - Inches(0.6), Inches(0.4), lbl,
                font=SANS_BOLD, size=12, bold=True, color=INK,
                char_spacing=80)
        textbox(s, x + Inches(0.3), base_y + Inches(2.0),
                cw - Inches(0.6), Inches(0.95), desc,
                font=SANS, size=10, color=INK2, line_spacing=1.4)
    textbox(s, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.3),
            "Sources: internal validation log · April 2026 · "
            "n=25 hardest ELEC3120 questions. Re-runnable live.",
            font="Cambria", size=10, italic=True, color=MUTED,
            align=PP_ALIGN.CENTER)
    return s


def slide_versus(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "07  ·  Built different")
    big_title(s, "side by side",
              "Same question. Two very different answers.")

    pw = (Inches(12.23) - Inches(0.4)) / 2
    px = [Inches(0.55), Inches(0.55) + pw + Inches(0.4)]
    py = Inches(3.4); ph = Inches(3.55)

    # ── ChatGPT panel ──
    cgpt = card(s, px[0], py, pw, ph, fill_c=GREY, stroke_c=HAIR,
                stroke_w=0.45, blur=16, dist=4, alpha=10, radius=0.04)
    textbox(s, px[0] + Inches(0.32), py + Inches(0.28),
            pw - Inches(0.6), Inches(0.4),
            "ChatGPT (free · no RAG)",
            font=SANS_BOLD, size=14, bold=True, color=INK)
    textbox(s, px[0] + Inches(0.32), py + Inches(0.78),
            pw - Inches(0.6), Inches(0.4),
            "Q.  In Reno, by what factor is cwnd reduced "
            "on a triple-dup ACK?",
            font="Cambria", size=11, italic=True, color=INK2)
    chat_bubble(s, px[0] + Inches(0.32), py + Inches(1.45),
                pw - Inches(0.64),
                text=("Reno reduces cwnd by ⅓ on a triple-duplicate ACK and "
                      "enters slow-start to recover quickly."),
                cite=None, h=Inches(1.2), font_size=11)
    rich(s, px[0] + Inches(0.32), py + Inches(2.85),
         pw - Inches(0.64), Inches(0.6),
         [[("✗  ", {"bold": True, "color": RED, "size": 13}),
           ("Wrong factor (½, not ⅓).",
            {"color": INK, "size": 11})],
          [("✗  ", {"bold": True, "color": RED, "size": 13}),
           ("Wrong recovery mode (cong-avoid, not slow-start).",
            {"color": INK, "size": 11})]],
         line_spacing=1.4)

    # ── LearningPacer panel ──
    lp = card(s, px[1], py, pw, ph, fill_c=PAPER, stroke_c=TEAL,
              stroke_w=0.9, blur=22, dist=6, alpha=16, radius=0.04)
    rect(s, px[1], py, pw, Inches(0.08), fill_c=TEAL,
         rounded=True, radius=0.5)
    textbox(s, px[1] + Inches(0.32), py + Inches(0.28),
            pw - Inches(0.6), Inches(0.4),
            "LearningPacer",
            font=SANS_BOLD, size=14, bold=True, color=TEAL)
    textbox(s, px[1] + Inches(0.32), py + Inches(0.78),
            pw - Inches(0.6), Inches(0.4),
            "Q.  In Reno, by what factor is cwnd reduced "
            "on a triple-dup ACK?",
            font="Cambria", size=11, italic=True, color=INK2)
    chat_bubble(s, px[1] + Inches(0.32), py + Inches(1.45),
                pw - Inches(0.64),
                text=("Reno halves cwnd (multiplicative decrease, factor ½) "
                      "and re-enters congestion-avoidance — not slow-start."),
                cite="Lec 09 · slide 17", h=Inches(1.2), font_size=11)
    rich(s, px[1] + Inches(0.32), py + Inches(2.85),
         pw - Inches(0.64), Inches(0.6),
         [[("✓  ", {"bold": True, "color": TEAL, "size": 13}),
           ("Right factor (½), right mode, right slide.",
            {"color": INK, "size": 11})],
          [("✓  ", {"bold": True, "color": TEAL, "size": 13}),
           ("Citation pill the student can verify in one click.",
            {"color": INK, "size": 11})]],
         line_spacing=1.4)

    textbox(s, Inches(0.55), Inches(7.05) - Inches(0.35),
            Inches(12.23), Inches(0.3),
            "ChatGPT is a brilliant generalist. LearningPacer knows ELEC3120.",
            font="Cambria", size=11, italic=True, color=MUTED,
            align=PP_ALIGN.CENTER)
    return s


def slide_demo_plan(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "08  ·  Live demo plan")
    big_title(s, "next eight minutes",
              "Six beats, in this order.")

    items = [
        ("00:30", "OPEN",     "Live app on conference wifi — no localhost, no screenshots."),
        ("02:00", "ASK",      "A hard ELEC3120 question (Reno vs Tahoe vs CUBIC, with derivation)."),
        ("03:30", "GENERATE", "A brand-new mock midterm in front of you, complete with diagrams."),
        ("05:30", "QUIZ",     "Five questions on a topic you pick — watch the stats ring update."),
        ("07:00", "BROWSE",   "Knowledge base · subnet calc · glossary in EN / 中文."),
        ("08:00", "Q & A",    "25-question pack rehearsed — but unrehearsed questions welcome."),
    ]
    base_x = Inches(0.7); base_y = Inches(3.3)
    row_h = Inches(0.6)
    rail_x = base_x + Inches(1.25)

    line(s, rail_x, base_y + Inches(0.2),
         rail_x, base_y + (len(items) - 0.65) * row_h,
         color=HAIR, w=1.2)

    for i, (t, head, body) in enumerate(items):
        ry = base_y + i * row_h
        textbox(s, base_x, ry + Inches(0.05),
                Inches(1.1), Inches(0.4), t,
                font=MONO, size=11, bold=True, color=TEAL,
                align=PP_ALIGN.RIGHT)
        # halo + dot
        oval(s, rail_x, ry + Inches(0.25), Inches(0.13),
             fill_c=TEAL_L)
        d = oval(s, rail_x, ry + Inches(0.25), Inches(0.085),
                 fill_c=TEAL, stroke_c=PAPER, stroke_w=1.5)
        add_shadow(d, blur=4, dist=1, alpha=22)
        textbox(s, rail_x + Inches(0.35), ry,
                Inches(2.6), Inches(0.4), head,
                font=SANS_BOLD, size=12.5, bold=True, color=INK,
                line_spacing=1.0, char_spacing=120)
        textbox(s, rail_x + Inches(2.95), ry + Inches(0.04),
                Inches(8.5), Inches(0.5), body,
                font=SANS, size=11, color=INK2, line_spacing=1.3)
    return s


def slide_roadmap(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "09  ·  Past the demo")
    big_title(s, "what's next",
              "Four steps from one course to a department deployment.")

    base_x = Inches(0.55); base_y = Inches(3.3); ch = Inches(3.55)
    span_w = Inches(12.23); n = 4; gap = Inches(0.18)
    cw = (span_w - (n - 1) * gap) / n
    rows = [
        ("MAY 2026", "Today",
         "FYP demo build · ELEC3120 only · single corpus."),
        ("Q3 2026",  "Multi-course",
         "ELEC3100, COMP3511 · per-course corpus flags · admin panel."),
        ("Q4 2026",  "Spaced repetition",
         "Weak-topic flashcards re-asked on a real schedule · cram pack."),
        ("2027",     "Department",
         "LTI · Canvas integration · anonymised analytics for staff."),
    ]
    for i, (when, head, body) in enumerate(rows):
        x = base_x + i * (cw + gap)
        c = card(s, x, base_y, cw, ch, fill_c=PAPER,
                 blur=18, dist=5, alpha=12)
        # teal date band
        rect(s, x, base_y, cw, Inches(0.62), fill_c=TEAL,
             rounded=True, radius=0.07)
        rect(s, x, base_y + Inches(0.32), cw, Inches(0.30), fill_c=TEAL)
        textbox(s, x + Inches(0.22), base_y + Inches(0.13),
                cw - Inches(0.44), Inches(0.4), when,
                font=MONO, size=11, bold=True, color=PAPER,
                anchor=MSO_ANCHOR.MIDDLE, char_spacing=180)
        textbox(s, x + Inches(0.22), base_y + Inches(0.85),
                cw - Inches(0.44), Inches(0.55), head,
                font=SERIF_BOLD, size=22, bold=True, color=INK,
                line_spacing=1.0)
        line(s, x + Inches(0.22), base_y + Inches(1.5),
             x + Inches(1.0), base_y + Inches(1.5),
             color=TEAL, w=1.8)
        textbox(s, x + Inches(0.22), base_y + Inches(1.7),
                cw - Inches(0.44), Inches(1.85), body,
                font=SANS, size=11, color=INK, line_spacing=1.45)
        if i < n - 1:
            arrow(s, x + cw + Inches(0.025),
                  base_y + Inches(0.31),
                  x + cw + gap - Inches(0.025),
                  base_y + Inches(0.31),
                  color=TEAL, w=1.4)
    return s


def slide_try(page, total):
    s = prs.slides.add_slide(BLANK)
    chrome(s, page, total, "10  ·  Try it now")
    big_title(s, "open it on your phone",
              "Live for the duration of this demo.")

    # left: stylised QR-tile
    qr_x = Inches(0.85); qr_y = Inches(3.2)
    qr_size = Inches(3.6)
    outer = card(s, qr_x - Inches(0.22), qr_y - Inches(0.22),
                 qr_size + Inches(0.44), qr_size + Inches(0.44),
                 fill_c=WHITE, stroke_c=HAIR, stroke_w=0.45,
                 blur=22, dist=6, alpha=18, radius=0.03)
    cell = qr_size / 7
    pattern = [
        "1110111",
        "1010101",
        "1110011",
        "0011010",
        "1100111",
        "1010101",
        "1110111",
    ]
    # draw pattern cells with subtle rounding for refinement
    for r_i, row in enumerate(pattern):
        for c_i, ch_ in enumerate(row):
            if ch_ == "1":
                rect(s, qr_x + c_i * cell, qr_y + r_i * cell,
                     cell, cell, fill_c=INK,
                     rounded=True, radius=0.18)
    # corner positioning eyes (rounded, layered)
    for (rx, ry) in [(qr_x, qr_y),
                     (qr_x + 4 * cell, qr_y),
                     (qr_x, qr_y + 4 * cell)]:
        rect(s, rx, ry, cell * 3, cell * 3, fill_c=TEAL,
             rounded=True, radius=0.25)
        rect(s, rx + cell * 0.5, ry + cell * 0.5,
             cell * 2, cell * 2, fill_c=PAPER,
             rounded=True, radius=0.25)
        rect(s, rx + cell, ry + cell, cell, cell, fill_c=TEAL,
             rounded=True, radius=0.25)

    # right column
    rx = Inches(5.2); ry = Inches(3.3)
    textbox(s, rx, ry, Inches(7.6), Inches(0.4),
            "OR TYPE",
            font=SANS_BOLD, size=10, bold=True, color=TEAL,
            char_spacing=200)
    pill = rect(s, rx, ry + Inches(0.45), Inches(7.6), Inches(0.95),
                fill_c=PAPER, stroke_c=TEAL, stroke_w=1.0,
                rounded=True, radius=0.18)
    add_shadow(pill, blur=14, dist=3, alpha=12)
    textbox(s, rx + Inches(0.3), ry + Inches(0.45),
            Inches(7.0), Inches(0.95),
            "learningpacer.app",
            font=MONO, size=28, bold=True, color=INK,
            anchor=MSO_ANCHOR.MIDDLE)

    # placeholder reminder — slot in the gap between URL pill (~4.7)
    # and the meta block (5.3). Well clear of the footer.
    textbox(s, rx, Inches(4.85), Inches(7.6), Inches(0.3),
            "(replace QR with a real one before printing)",
            font="Cambria", size=9, italic=True, color=MUTED)

    ty = ry + Inches(2.0)
    items = [
        ("BUILT BY",   "[Your Name]"),
        ("SUPERVISOR", "[Supervisor]"),
        ("COURSE",     "ELEC3120 Computer Networks"),
        ("CORPUS",     "17 lecture decks · ~340 slides"),
    ]
    for i, (k, v) in enumerate(items):
        ly = ty + i * Inches(0.42)
        textbox(s, rx, ly, Inches(2.4), Inches(0.4), k,
                font=SANS_BOLD, size=10, bold=True, color=TEAL,
                char_spacing=180)
        textbox(s, rx + Inches(2.5), ly, Inches(5.0), Inches(0.4), v,
                font=SANS, size=12, color=INK)
    return s


def slide_close(page, total):
    s = prs.slides.add_slide(BLANK)
    paper(s)

    band = rect(s, 0, Inches(5.2), SW, SH - Inches(5.2), fill_c=TEAL)
    _no_outline(band)

    textbox(s, Inches(0.7), Inches(0.85), Inches(8), Inches(0.4),
            "11  ·  THANK YOU",
            font=SANS_BOLD, size=10, bold=True, color=TEAL,
            char_spacing=200)
    textbox(s, Inches(0.7), Inches(1.55), Inches(12), Inches(3.0),
            "Questions are how\nthe tutor learns, too.",
            font=SERIF_BOLD, size=72, bold=True, color=INK,
            line_spacing=1.05)
    line(s, Inches(0.7), Inches(4.85), Inches(5.5), Inches(4.85),
         color=TEAL, w=2.0)

    rich(s, Inches(0.7), Inches(5.55), Inches(12.0), Inches(0.5),
         [[("OPEN TO   ",
            {"bold": True, "color": PAPER, "size": 11}),
           ("code review · panel questions · unrehearsed prompts · "
            "live retries against ChatGPT, Gemini, anything.",
            {"color": PAPER, "size": 13})]],
         line_spacing=1.4)
    rich(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.5),
         [[("URL       ", {"bold": True, "color": PAPER, "size": 11,
                            "font": MONO}),
           ("learningpacer.app",
            {"font": MONO, "color": PAPER, "size": 13})]],
         line_spacing=1.4)
    textbox(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.3),
            "HKUST · Department of Electronic & Computer Engineering · "
            "Final-Year Project · May 2026",
            font=SANS, size=9, color=PAPER)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# build
# ─────────────────────────────────────────────────────────────────────────────
def build():
    BUILDERS = [
        ("cover", slide_cover),
        slide_moment,
        slide_meet,
        slide_anatomy,
        slide_pipeline,
        slide_moments,
        slide_numbers,
        slide_versus,
        slide_demo_plan,
        slide_roadmap,
        slide_try,
        slide_close,
    ]
    total = len(BUILDERS)
    page = 1
    for entry in BUILDERS:
        if isinstance(entry, tuple) and entry[0] == "cover":
            entry[1]()
        else:
            entry(page, total)
        page += 1
    prs.save(OUT_PATH)
    return OUT_PATH, total


if __name__ == "__main__":
    path, n = build()
    kb = os.path.getsize(path) / 1024
    print(f"OK  {path}  ({kb:.1f} KB, {n} slides)")
