"""LearningPacer FYP demo deck — PPTX edition.
16:9 widescreen · Claude warm-paper aesthetic · native editable shapes
Output: .local/fyp-deck/LearningPacer_FYP_Demo_Deck.pptx
Run:    python3 scripts/deck-gen/build_fyp_pptx.py
"""
import os, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── palette ─────────────────────────────────────────────────────────────────
PAPER  = RGBColor(0xfa, 0xf6, 0xee)   # warm cream bg
INK    = RGBColor(0x1e, 0x29, 0x3b)   # deep slate
INK2   = RGBColor(0x47, 0x55, 0x69)   # secondary
MUTED  = RGBColor(0x94, 0xa3, 0xb8)   # caption
HAIR   = RGBColor(0xcb, 0xd5, 0xe1)   # hairline
TEAL   = RGBColor(0x0f, 0x76, 0x6e)   # primary
CYAN   = RGBColor(0x08, 0x91, 0xb2)   # secondary
AMBER  = RGBColor(0xb4, 0x53, 0x09)   # warn
RED    = RGBColor(0xb9, 0x1c, 0x1c)   # danger
SOFT   = RGBColor(0xef, 0xe9, 0xda)   # tinted card
SOFT2  = RGBColor(0xe6, 0xdf, 0xcc)
WHITE  = RGBColor(0xff, 0xff, 0xff)

# ── font choices (PowerPoint-safe) ───────────────────────────────────────────
SERIF_BOLD = "Cambria"          # headline serif (every Mac/Windows ship it)
SANS       = "Calibri"          # body
SANS_BOLD  = "Calibri"
MONO       = "Consolas"

# ── canvas: 16:9 widescreen ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

OUT_DIR  = ".local/fyp-deck"
os.makedirs(OUT_DIR, exist_ok=True)
OUT_PATH = os.path.join(OUT_DIR, "LearningPacer_FYP_Demo_Deck.pptx")

BLANK = prs.slide_layouts[6]   # fully blank


# ── primitives ───────────────────────────────────────────────────────────────
def _no_outline(shape):
    shape.line.fill.background()


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def stroke(shape, color, width_pt=0.6):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def rect(slide, x, y, w, h, fill_c=None, stroke_c=None, stroke_w=0.6, rounded=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, x, y, w, h)
    s.shadow.inherit = False
    if fill_c is not None:
        fill(s, fill_c)
    else:
        s.fill.background()
    if stroke_c is not None:
        stroke(s, stroke_c, stroke_w)
    else:
        _no_outline(s)
    if rounded:
        # tighten corner radius
        try:
            s.adjustments[0] = 0.08
        except Exception:
            pass
    return s


def line(slide, x1, y1, x2, y2, color=HAIR, width_pt=0.4, dash=False):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    s.line.color.rgb = color
    s.line.width = Pt(width_pt)
    if dash:
        # apply prstDash="dash" via XML
        ln = s.line._get_or_add_ln()
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return s


def arrow(slide, x1, y1, x2, y2, color=TEAL, width_pt=0.9, dash=False):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    s.line.color.rgb = color
    s.line.width = Pt(width_pt)
    ln = s.line._get_or_add_ln()
    if dash:
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return s


def textbox(slide, x, y, w, h, text, *,
            font=SANS, size=12, bold=False, color=INK,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rich_textbox(slide, x, y, w, h, runs, *,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs = list[list[(text, dict_of_style)]]   — one inner list per paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, st in para_runs:
            r = p.add_run()
            r.text = txt
            r.font.name = st.get("font", SANS)
            r.font.size = Pt(st.get("size", 12))
            r.font.bold = st.get("bold", False)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", INK)
    return tb


# ── chrome ───────────────────────────────────────────────────────────────────
def add_paper(slide):
    bg = rect(slide, 0, 0, SW, SH, fill_c=PAPER)
    bg.line.fill.background()
    return bg


def add_chrome_normal(slide, page_num, total, eyebrow_left=None):
    add_paper(slide)
    # full-height left teal sidebar (4mm-ish = ~Inches(0.18))
    sb = rect(slide, 0, 0, Inches(0.18), SH, fill_c=TEAL)
    _no_outline(sb)
    # top eyebrow
    textbox(slide, Inches(0.55), Inches(0.28), Inches(7), Inches(0.3),
            "LEARNINGPACER  ·  ELEC3120 VIRTUAL TA",
            font=SANS_BOLD, size=8, bold=True, color=TEAL)
    textbox(slide, Inches(7.2),  Inches(0.28), Inches(5.6), Inches(0.3),
            "HKUST  ·  Final-Year Project  ·  2025–2026",
            font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
    # top hairline
    line(slide, Inches(0.55), Inches(0.62), Inches(12.78), Inches(0.62),
         color=HAIR, width_pt=0.4)
    # bottom hairline + footer
    line(slide, Inches(0.55), Inches(7.05), Inches(12.78), Inches(7.05),
         color=HAIR, width_pt=0.4)
    textbox(slide, Inches(0.55), Inches(7.15), Inches(8), Inches(0.3),
            "Confidential demo material  ·  Dept of ECE",
            font=SANS, size=8, color=MUTED)
    textbox(slide, Inches(8.5), Inches(7.15), Inches(4.3), Inches(0.3),
            f"{page_num:02d} / {total:02d}",
            font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_header(slide, eyebrow, title, sub=None, y_top=Inches(0.95)):
    """Reusable interior-slide header block."""
    textbox(slide, Inches(0.55), y_top, Inches(12), Inches(0.32),
            eyebrow.upper(), font=SANS_BOLD, size=9, bold=True, color=TEAL)
    textbox(slide, Inches(0.55), y_top + Inches(0.3), Inches(12), Inches(0.7),
            title, font=SERIF_BOLD, size=30, bold=True, color=INK,
            line_spacing=1.05)
    if sub:
        textbox(slide, Inches(0.55), y_top + Inches(1.0), Inches(12),
                Inches(0.45), sub, font=SANS, size=14, color=INK2,
                line_spacing=1.25)
    # divider rule
    rule_y = y_top + Inches(1.45 if sub else 1.0)
    line(slide, Inches(0.55), rule_y, Inches(12.78), rule_y,
         color=HAIR, width_pt=0.4)
    return rule_y + Inches(0.18)   # next-y where content can start


# ── reusable shape: stat card ────────────────────────────────────────────────
def stat_column(slide, x, y, w, h, num, label, desc):
    # auto-shrink long stat numbers
    sz = 64 if len(num) <= 4 else 40
    textbox(slide, x, y, w, Inches(1.1), num,
            font=SERIF_BOLD, size=sz, bold=True, color=TEAL,
            line_spacing=1.0)
    textbox(slide, x, y + Inches(1.15), w, Inches(0.4), label,
            font=SANS_BOLD, size=12, bold=True, color=INK)
    textbox(slide, x, y + Inches(1.55), w, Inches(2.0), desc,
            font=SANS, size=10, color=INK2, line_spacing=1.35)


# ── concentric brand mark ────────────────────────────────────────────────────
def brand_mark(slide, cx, cy, max_r):
    """Concentric ripple in teal — cover-page lockup."""
    radii = [max_r, max_r*0.82, max_r*0.62, max_r*0.43, max_r*0.26, max_r*0.13]
    for r in radii:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   cx - r, cy - r, r * 2, r * 2)
        s.shadow.inherit = False
        s.fill.background()
        s.line.color.rgb = TEAL
        s.line.width = Pt(0.9)
    # core dot
    core_r = max_r * 0.05
    core = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  cx - core_r, cy - core_r,
                                  core_r * 2, core_r * 2)
    core.shadow.inherit = False
    fill(core, TEAL); _no_outline(core)
    # cardinal ticks
    tick = max_r * 0.06
    for ang in (0, 90, 180, 270):
        rad = math.radians(ang)
        x1 = cx + (max_r) * math.cos(rad)
        y1 = cy + (max_r) * math.sin(rad)
        x2 = cx + (max_r + tick) * math.cos(rad)
        y2 = cy + (max_r + tick) * math.sin(rad)
        line(slide, x1, y1, x2, y2, color=TEAL, width_pt=1.2)


# ── architecture diagram ────────────────────────────────────────────────────
def arch_diagram(slide, x, y, w, h):
    """Native PPTX architecture diagram — fully editable."""
    # background card (subtle)
    bg = rect(slide, x, y, w, h, fill_c=PAPER, stroke_c=HAIR, stroke_w=0.4)

    title = textbox(slide, x + Inches(0.2), y + Inches(0.12),
                    w - Inches(0.4), Inches(0.3),
                    "System architecture — request flow",
                    font=SANS_BOLD, size=10, bold=True, color=TEAL)

    # node helper
    def node(nx, ny, nw, nh, label, sub=None,
             fill_c=SOFT, stroke_c=TEAL):
        s = rect(slide, nx, ny, nw, nh, fill_c=fill_c, stroke_c=stroke_c,
                 stroke_w=0.7, rounded=True)
        if sub:
            tb = slide.shapes.add_textbox(nx, ny, nw, nh)
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run(); r1.text = label
            r1.font.name = SANS_BOLD; r1.font.size = Pt(11); r1.font.bold = True
            r1.font.color.rgb = INK
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            p2.line_spacing = 1.15
            r2 = p2.add_run(); r2.text = sub
            r2.font.name = SANS; r2.font.size = Pt(8); r2.font.color.rgb = INK2
        else:
            tb = slide.shapes.add_textbox(nx, ny, nw, nh)
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = Emu(0)
            tf.margin_top = tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = label
            r.font.name = SANS_BOLD; r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = INK

    # Diagram coordinate frame: inset by 0.4"
    DX = x + Inches(0.4); DY = y + Inches(0.55)
    DW = w - Inches(0.8); DH = h - Inches(0.7)

    # Two rows of nodes; column x positions
    cw  = Inches(1.85); ch_top = Inches(0.7); ch_bot = Inches(0.7)
    gap = Inches(0.25)
    col_x = [DX + i * (cw + gap) for i in range(5)]

    # Top row (request chain)
    top_y = DY + Inches(0.05)
    node(col_x[0], top_y, cw, ch_top, "Browser",      "Next.js 16 / React 19")
    node(col_x[1], top_y, cw, ch_top, "Edge API",     "Route Handlers")
    node(col_x[2], top_y, cw, ch_top, "Tutor Engine", "prompt + RAG + tools")
    node(col_x[3], top_y, cw, Inches(0.42), "Anthropic",  "claude-sonnet-4")
    node(col_x[3], top_y + Inches(0.5), cw, Inches(0.42), "OpenRouter",
         "fallback fleet")

    # Bottom row (data)
    bot_y = DY + Inches(1.85)
    node(col_x[0], bot_y, cw, ch_bot, "Supabase",   "auth · Postgres · RLS")
    node(col_x[1], bot_y, cw, ch_bot, "RAG Index",  "lecture PDFs · vectors")
    node(col_x[2], bot_y, cw, ch_bot, "Tool Belt",  "subnet · diagram · search")
    node(col_x[3], bot_y, cw, ch_bot, "Storage",    "user uploads")

    # Arrows — request chain
    def cy(yy, cc): return yy + cc / 2
    arrow(slide, col_x[0] + cw, cy(top_y, ch_top),
                 col_x[1],       cy(top_y, ch_top), color=TEAL)
    arrow(slide, col_x[1] + cw, cy(top_y, ch_top),
                 col_x[2],       cy(top_y, ch_top), color=TEAL)
    arrow(slide, col_x[2] + cw, cy(top_y, ch_top),
                 col_x[3],       top_y + Inches(0.21), color=TEAL)
    arrow(slide, col_x[2] + cw, cy(top_y, ch_top),
                 col_x[3],       top_y + Inches(0.71), color=TEAL)
    # response back (dashed cyan)
    arrow(slide, col_x[1], cy(top_y, ch_top) + Inches(0.18),
                 col_x[0] + cw, cy(top_y, ch_top) + Inches(0.18),
                 color=CYAN, dash=True, width_pt=0.7)
    # vertical: api → data
    arrow(slide, col_x[0] + cw/2, top_y + ch_top,
                 col_x[0] + cw/2, bot_y, color=TEAL, width_pt=0.7)
    arrow(slide, col_x[1] + cw/2, top_y + ch_top,
                 col_x[1] + cw/2, bot_y, color=TEAL, width_pt=0.7)
    arrow(slide, col_x[2] + cw/2, top_y + ch_top,
                 col_x[2] + cw/2, bot_y, color=TEAL, width_pt=0.7)

    # Feature banner across bottom
    fb_y = DY + DH - Inches(0.55)
    fb_w = (col_x[3] + cw) - col_x[0]
    fb = rect(slide, col_x[0], fb_y, fb_w, Inches(0.5),
              fill_c=PAPER, stroke_c=CYAN, stroke_w=0.7, rounded=True)
    rich_textbox(slide, col_x[0], fb_y, fb_w, Inches(0.5), [[
        ("Features delivered  ", {"bold": True, "size": 10, "color": INK}),
        ("Tutor · Mock exams · Quiz · Pomodoro · Stats · Flashcards · "
         "Glossary · Subnet · Net tools",
         {"size": 9, "color": INK2}),
    ]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── three-pillar cards ──────────────────────────────────────────────────────
def three_pillars(slide, x, y, w, h):
    cards = [
        ("01", "Conversational Tutor",
         "Citation-grounded answers in EN/中文.",
         "Every reply traces back to a lecture slide."),
        ("02", "Mock Exam Engine",
         "Full midterm + final papers on demand.",
         "Worked solutions, diagrams, marks breakdown."),
        ("03", "Personal Study Layer",
         "Quiz · Pomodoro · Stats · Goals · Notes.",
         "Tracks every minute; surfaces weak topics."),
    ]
    gap = Inches(0.25)
    cw = (w - 2 * gap) / 3
    for i, (n, t, l1, l2) in enumerate(cards):
        cx = x + i * (cw + gap)
        # card body
        body = rect(slide, cx, y, cw, h,
                    fill_c=SOFT, stroke_c=TEAL, stroke_w=0.6, rounded=True)
        # teal top stripe
        stripe = rect(slide, cx, y, cw, Inches(0.55), fill_c=TEAL, rounded=True)
        # square the bottom of the stripe by overlay rect
        rect(slide, cx, y + Inches(0.30), cw, Inches(0.25), fill_c=TEAL)
        # number
        textbox(slide, cx + Inches(0.18), y + Inches(0.08),
                Inches(0.5), Inches(0.4), n,
                font=SERIF_BOLD, size=14, bold=True, color=WHITE)
        # title
        textbox(slide, cx + Inches(0.18), y + Inches(0.7),
                cw - Inches(0.36), Inches(0.5), t,
                font=SANS_BOLD, size=15, bold=True, color=INK)
        # line 1 + line 2
        textbox(slide, cx + Inches(0.18), y + Inches(1.18),
                cw - Inches(0.36), Inches(0.7), l1,
                font=SANS, size=11, color=INK2, line_spacing=1.3)
        textbox(slide, cx + Inches(0.18), y + Inches(1.78),
                cw - Inches(0.36), Inches(1.0), l2,
                font=SANS, size=10, color=MUTED, line_spacing=1.3)


# ── demo timeline ───────────────────────────────────────────────────────────
def demo_timeline(slide, x, y, w, h):
    steps = [
        ("00:30", "Launch"),
        ("02:00", "Hard Q\n+ RAG"),
        ("03:30", "Mock\npaper"),
        ("05:30", "Quiz +\nPomodoro"),
        ("07:00", "Stats +\nK-base"),
        ("08:00", "Q&A"),
    ]
    n = len(steps)
    pad = Inches(0.6)
    rail_y = y + h / 2
    span = w - 2 * pad

    # rail
    line(slide, x + pad, rail_y, x + w - pad, rail_y,
         color=HAIR, width_pt=1.0)
    line(slide, x + pad, rail_y, x + w - pad, rail_y,
         color=TEAL, width_pt=1.0, dash=True)

    dot_r = Inches(0.09)
    for i, (t, lbl) in enumerate(steps):
        cx = x + pad + i * span / (n - 1)
        # dot
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    cx - dot_r, rail_y - dot_r,
                                    dot_r * 2, dot_r * 2)
        d.shadow.inherit = False
        fill(d, TEAL); stroke(d, PAPER, 1.5)
        # time
        textbox(slide, cx - Inches(0.6), rail_y - Inches(0.55),
                Inches(1.2), Inches(0.3), t,
                font=SANS_BOLD, size=10, bold=True, color=INK,
                align=PP_ALIGN.CENTER)
        # label
        textbox(slide, cx - Inches(0.7), rail_y + Inches(0.12),
                Inches(1.4), Inches(0.7), lbl,
                font=SANS, size=10, color=INK2, align=PP_ALIGN.CENTER,
                line_spacing=1.2)


# ── slide builders ───────────────────────────────────────────────────────────
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_paper(s)
    # full-height teal sidebar — wide for cover (1.3")
    sb = rect(s, 0, 0, Inches(1.3), SH, fill_c=TEAL); _no_outline(sb)
    # left column text
    textbox(s, Inches(1.65), Inches(0.85), Inches(8), Inches(0.4),
            "LEARNINGPACER", font=SANS_BOLD, size=11, bold=True, color=TEAL)
    # accent rule under brand
    line(s, Inches(1.65), Inches(1.25), Inches(6.8), Inches(1.25),
         color=TEAL, width_pt=1.5)
    # hero title (two lines)
    textbox(s, Inches(1.65), Inches(1.6), Inches(8.2), Inches(2.6),
            "The 24/7 virtual TA\nfor ELEC3120.",
            font=SERIF_BOLD, size=58, bold=True, color=INK,
            line_spacing=1.05)
    textbox(s, Inches(1.65), Inches(4.2), Inches(8.2), Inches(0.9),
            "A grounded, exam-grade study companion for "
            "Computer Networks at HKUST.",
            font="Cambria", size=18, color=INK2, line_spacing=1.35)
    # meta block (right under sub)
    rich_textbox(s, Inches(1.65), Inches(5.3), Inches(8.2), Inches(1.6), [
        [("Final-Year Project Demo  ·  May 2026",
          {"size": 11, "color": INK2})],
        [("Presenter:   ", {"bold": True, "size": 11, "color": INK}),
         ("[Your Name]",   {"size": 11, "color": INK2})],
        [("Supervisor:  ", {"bold": True, "size": 11, "color": INK}),
         ("[Supervisor]",  {"size": 11, "color": INK2})],
        [("Course:      ", {"bold": True, "size": 11, "color": INK}),
         ("ELEC3120 Computer Networks", {"size": 11, "color": INK2})],
    ], line_spacing=1.5)

    # right brand mark
    brand_mark(s, Inches(11.0), Inches(3.7), Inches(2.2))

    # bottom rule + footer (cover)
    line(s, Inches(1.65), Inches(7.05), Inches(12.78), Inches(7.05),
         color=HAIR, width_pt=0.4)
    textbox(s, Inches(1.65), Inches(7.15), Inches(8.5), Inches(0.3),
            "Hong Kong University of Science and Technology  ·  "
            "Department of Electronic & Computer Engineering",
            font=SANS, size=8, color=MUTED)
    textbox(s, Inches(10.5), Inches(7.15), Inches(2.3), Inches(0.3),
            "FYP DEMO  ·  v2.0",
            font=SANS, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def slide_context(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "01 · Context", "ELEC3120, in numbers.",
               "Where the pressure on a single course actually lives.")
    stats = [
        ("~200",    "Students per cohort",
         "One lecturer; no dedicated TA office hours stream."),
        ("17",      "Lecture decks",
         "OSI to wireless — over 300 dense slides of material."),
        ("14 days", "Until midterm + final",
         "Both papers compress into the same revision window each term."),
        ("∞",       "Questions at 2 a.m.",
         "Right before the exam, students have nowhere to ask."),
    ]
    base_x = Inches(0.55); base_y = Inches(3.2)
    col_w = Inches(3.0); gap = Inches(0.15)
    for i, (n, l, d) in enumerate(stats):
        cx = base_x + i * (col_w + gap)
        if i > 0:
            line(s, cx - gap/2, base_y, cx - gap/2, base_y + Inches(3.5),
                 color=HAIR, width_pt=0.4)
        stat_column(s, cx + Inches(0.1), base_y, col_w - Inches(0.2),
                    Inches(3.5), n, l, d)
    return s


def slide_problem(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "02 · The problem",
               "A content problem that is really a dialogue problem.",
               "The lecture material exists. Grounded, on-demand dialogue does not.")
    # divider between two columns
    line(s, Inches(6.65), Inches(3.2), Inches(6.65), Inches(6.85),
         color=HAIR, width_pt=0.4)

    # left
    textbox(s, Inches(0.55), Inches(3.15), Inches(5.8), Inches(0.3),
            "WHAT BREAKS DOWN", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    bullets_l = [
        "Office hours fill up days before the midterm. When students need answers, the queue is closed.",
        "Generic chatbots hallucinate confidently on protocols, subnetting, and ELEC3120 exam conventions.",
        "Past papers, model answers, and trap patterns are scattered across PDFs, WhatsApp groups, and seniors' Drives.",
        "Self-study tools live in separate apps that don't know what you got wrong last Tuesday.",
    ]
    rich_textbox(s, Inches(0.55), Inches(3.5), Inches(5.8), Inches(3.3),
        [[("•  " + b, {"size": 11, "color": INK})] for b in bullets_l],
        line_spacing=1.4)

    # right
    textbox(s, Inches(6.95), Inches(3.15), Inches(5.85), Inches(0.3),
            "WHAT STUDENTS ACTUALLY DO",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(6.95), Inches(3.5), Inches(5.85), Inches(0.95),
            "They paste lecture slides into ChatGPT, get plausible-sounding "
            "but subtly wrong answers, and ship them straight into revision notes.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    # quote block
    qbg = rect(s, Inches(6.95), Inches(4.55), Inches(5.85), Inches(1.05),
               fill_c=SOFT, rounded=True)
    qbg.line.fill.background()
    rich_textbox(s, Inches(7.15), Inches(4.65), Inches(5.55), Inches(0.95), [
        [("\u201cIt sounded right, so I memorised it.\u201d",
          {"font": "Cambria", "size": 14, "italic": True, "color": INK})],
        [("— ELEC3120 student, two days before the final",
          {"size": 9, "color": MUTED})],
    ], line_spacing=1.4)
    textbox(s, Inches(6.95), Inches(5.85), Inches(5.85), Inches(1.0),
            "The failure mode is not laziness. It is the absence of a "
            "trustworthy, always-on interlocutor who knows the course.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    return s


def slide_solution(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "03 · Solution", "One companion, three pillars.",
               "Tutoring, examination, and self-management — fused into one revision loop.")
    three_pillars(s, Inches(0.55), Inches(3.25), Inches(12.23), Inches(3.0))
    textbox(s, Inches(0.55), Inches(6.45), Inches(12.23), Inches(0.5),
            "Every pillar draws from the same lecture corpus. The model "
            "never free-styles when the right answer is already in the slides.",
            font=SANS, size=12, color=INK, line_spacing=1.4)
    return s


def slide_architecture(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "04 · Architecture",
               "A thin, opinionated stack — built for trust.",
               "Modern web primitives, two LLM providers, explicit retrieval grounding.")
    arch_diagram(s, Inches(0.55), Inches(3.05), Inches(12.23), Inches(3.55))
    textbox(s, Inches(0.55), Inches(6.7), Inches(12.23), Inches(0.35),
            "Anthropic Claude is primary; OpenRouter provides a fallback fleet "
            "for cost ceilings and uptime. Supabase Row-Level Security ensures "
            "each student only ever reads their own data.",
            font=SANS, size=10, color=INK2, line_spacing=1.35)
    return s


def slide_stack(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "05 · Stack", "Twelve choices, one defensible reason each.")
    rows = [
        ("Frontend",
         "Next.js 16 (App Router) · React 19 · Tailwind v4 · shadcn/ui",
         "Streaming UI, server components, file-based routing — zero config."),
        ("Auth & DB",
         "Supabase · PostgreSQL · Row-Level Security",
         "Magic-link auth; per-row policies; zero custom auth code to audit."),
        ("Primary LLM",
         "Anthropic Claude Sonnet 4 · OpenRouter fallback fleet",
         "Best reasoning + reliable JSON tool use; second provider for cost & uptime."),
        ("Grounding",
         "Lecture PDF embeddings · chunk-and-cite RAG pipeline",
         "Every answer carries a slide-deck citation the student can verify."),
        ("Tool belt",
         "Subnet calc · diagram engine · web search · code runner",
         "Domain tools the model invokes when the question demands them."),
        ("Hosting",
         "Replit Deployments · autoscaling · shareable preview URLs",
         "One-click publish; shareable link for marking; scales to exam-week traffic."),
    ]
    base_y = Inches(2.85)
    row_h = Inches(0.62)
    col_w = (Inches(2.0), Inches(5.4), Inches(4.83))
    cx = [Inches(0.55), Inches(0.55) + col_w[0], Inches(0.55) + col_w[0] + col_w[1]]
    for i, (label, picks, why) in enumerate(rows):
        ry = base_y + i * row_h
        if i > 0:
            line(s, Inches(0.55), ry, Inches(12.78), ry,
                 color=HAIR, width_pt=0.3)
        textbox(s, cx[0], ry + Inches(0.08), col_w[0], row_h, label,
                font=SANS_BOLD, size=11, bold=True, color=INK)
        textbox(s, cx[1], ry + Inches(0.08), col_w[1], row_h, picks,
                font=SANS, size=11, color=INK)
        textbox(s, cx[2], ry + Inches(0.08), col_w[2], row_h, why,
                font=SANS, size=10, color=INK2, line_spacing=1.3)
    return s


def slide_tutor(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "06 · Feature — Tutor",
               "Grounded conversation, not autocomplete.",
               "Answering ELEC3120 questions feels like talking to the top student in the cohort.")

    line(s, Inches(6.65), Inches(3.2), Inches(6.65), Inches(6.85),
         color=HAIR, width_pt=0.4)

    textbox(s, Inches(0.55), Inches(3.15), Inches(5.8), Inches(0.3),
            "HOW IT DIFFERS", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    rich_textbox(s, Inches(0.55), Inches(3.5), Inches(5.8), Inches(2.3), [
        [("•  Every reply cites the exact lecture deck and slide number — verifiable, not vibes.",
          {"size": 11, "color": INK})],
        [("•  Stays inside the ELEC3120 ontology: Prof's notation, diagram style, marking conventions.",
          {"size": 11, "color": INK})],
        [("•  Switches EN / 中文 inline, matching how Cantonese students actually think during revision.",
          {"size": 11, "color": INK})],
        [("•  Refuses confidently when a question is out of scope — no plausible fabrications.",
          {"size": 11, "color": INK})],
    ], line_spacing=1.45)
    textbox(s, Inches(0.55), Inches(5.85), Inches(5.8), Inches(0.3),
            "FAILURE MODES WE HANDLE",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    rich_textbox(s, Inches(0.55), Inches(6.2), Inches(5.8), Inches(0.95), [
        [("•  Off-by-one errors in subnetting questions",
          {"size": 11, "color": INK})],
        [("•  Tahoe vs Reno vs CUBIC congestion-control terminology",
          {"size": 11, "color": INK})],
        [("•  Chinese-medium ambiguity on protocol acronyms",
          {"size": 11, "color": INK})],
    ], line_spacing=1.4)

    textbox(s, Inches(6.95), Inches(3.15), Inches(5.85), Inches(0.3),
            "UI MODES", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(6.95), Inches(3.5), Inches(5.85), Inches(0.4),
            "Tutor · Code · Image · Web search · Agent",
            font=SANS, size=11, color=INK)
    textbox(s, Inches(6.95), Inches(4.0), Inches(5.85), Inches(0.3),
            "QUICK-LAUNCH CHIPS", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(6.95), Inches(4.35), Inches(5.85), Inches(0.95),
            "Mock paper (with diagrams) · Lecture quiz · Key points · "
            "Plain explain · Key formulas · Exam tips",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    # quote
    qbg = rect(s, Inches(6.95), Inches(5.55), Inches(5.85), Inches(1.5),
               fill_c=SOFT, rounded=True); qbg.line.fill.background()
    rich_textbox(s, Inches(7.15), Inches(5.7), Inches(5.55), Inches(1.3), [
        [("\u201cExplain why TCP Reno backs off to half cwnd, with the slide reference and a small diagram.\u201d",
          {"font": "Cambria", "size": 13, "italic": True, "color": INK})],
        [("→ Cited from Lecture 07, slide 14. Diagram rendered inline.",
          {"size": 9.5, "color": INK2})],
    ], line_spacing=1.4)
    return s


def slide_exam(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "07 · Feature — Exam engine",
               "From a one-line prompt to a printable mock paper.",
               "Lecturer-style midterm + final, with worked solutions and rendered diagrams.")
    line(s, Inches(6.65), Inches(3.2), Inches(6.65), Inches(6.85),
         color=HAIR, width_pt=0.4)

    textbox(s, Inches(0.55), Inches(3.15), Inches(5.8), Inches(0.3),
            "WHAT WE GENERATE TODAY",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    rich_textbox(s, Inches(0.55), Inches(3.5), Inches(5.8), Inches(2.6), [
        [("•  Full mock midterm + final (33 pages, 8 sections, marks allocated, model solutions).",
          {"size": 11, "color": INK})],
        [("•  Trap-analysis booklet (14 pages of common-mistake patterns, each tied to a real exam question).",
          {"size": 11, "color": INK})],
        [("•  Native diagrams: OSI stack, TCP handshake, sliding window, congestion control, subnet, packet journey.",
          {"size": 11, "color": INK})],
        [("•  25-question panel-defence pack with written model answers.",
          {"size": 11, "color": INK})],
    ], line_spacing=1.4)
    textbox(s, Inches(0.55), Inches(6.15), Inches(5.8), Inches(0.3),
            "HOW IT STAYS HONEST",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    rich_textbox(s, Inches(0.55), Inches(6.5), Inches(5.8), Inches(0.7), [
        [("•  Every question is anchored to a specific lecture slide.",
          {"size": 11, "color": INK})],
        [("•  Difficulty distribution mirrors real ELEC3120 papers (definition · derivation · scenario · numerical).",
          {"size": 11, "color": INK})],
    ], line_spacing=1.4)

    textbox(s, Inches(6.95), Inches(3.15), Inches(5.85), Inches(0.3),
            "ARTIFACTS ALREADY ON DISK",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    rich_textbox(s, Inches(6.95), Inches(3.5), Inches(5.85), Inches(2.4), [
        [("ELEC3120_Mock_Midterm_and_Final.pdf",
          {"bold": True, "size": 11, "color": INK})],
        [("33 pp · 8 sections · model solutions",
          {"size": 10, "color": INK2})],
        [(" ", {"size": 4})],
        [("ELEC3120_Trap_Analysis_and_Answering_Techniques.pdf",
          {"bold": True, "size": 11, "color": INK})],
        [("14 pp · pattern library for the most-missed topics",
          {"size": 10, "color": INK2})],
        [(" ", {"size": 4})],
        [("FYP Defence Q&A Pack",
          {"bold": True, "size": 11, "color": INK})],
        [("25 panel questions with prepared model answers",
          {"size": 10, "color": INK2})],
    ], line_spacing=1.35)
    textbox(s, Inches(6.95), Inches(6.0), Inches(5.85), Inches(1.0),
            "All three export as print-ready PDFs. Professors can mark "
            "them on paper without opening a laptop.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    return s


def slide_study(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "08 · Feature — Personal study layer",
               "The boring 90% of revision, finally in one place.")
    items = [
        ("Quiz",       "Lecture-scoped MCQ; explanations after each answer."),
        ("Pomodoro",   "Inline timer; minutes count toward the daily goal ring."),
        ("Stats",      "Per-day study minutes, streaks, weak topics surfaced."),
        ("Plan",       "Two-week revision plan from exam date + weak topics."),
        ("Notes",      "Free-form notes the tutor can read before answering."),
        ("Flashcards", "Auto-generated from any chat; spaced-repetition queue."),
        ("Formulas",   "One canonical sheet — searchable, copyable."),
        ("Glossary",   "Per-protocol terms in EN and 中文."),
        ("Protocols",  "Comparative breakdown: TCP / UDP / IP / BGP / ARP / ..."),
        ("Subnet",     "Calculator with full worked steps, not just the answer."),
        ("Net Tools",  "ping · traceroute · DNS lookup, illustrated for teaching."),
        ("Goals",      "Daily minutes goal; visualised as a ring on the sidebar."),
    ]
    base_x = Inches(0.55); base_y = Inches(2.85)
    cw = Inches(3.0); ch = Inches(1.35); gap = Inches(0.08)
    for i, (name, desc) in enumerate(items):
        col = i % 4; row = i // 4
        cx = base_x + col * (cw + gap)
        cy = base_y + row * (ch + gap)
        card = rect(s, cx, cy, cw, ch, fill_c=SOFT, stroke_c=HAIR,
                    stroke_w=0.4, rounded=True)
        textbox(s, cx + Inches(0.18), cy + Inches(0.1),
                cw - Inches(0.3), Inches(0.35), name,
                font=SANS_BOLD, size=12, bold=True, color=TEAL)
        textbox(s, cx + Inches(0.18), cy + Inches(0.5),
                cw - Inches(0.3), Inches(0.85), desc,
                font=SANS, size=10, color=INK2, line_spacing=1.3)
    return s


def slide_pedagogy(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "09 · Pedagogy",
               "Why this works for ELEC3120 specifically.",
               "Teaching choices, not just engineering choices.")
    rows = [
        ("Worked-example bias",
         "Every explanation includes a derivation or diagram, never just a verdict.",
         "Mirrors how the lecturer teaches: derive first, then assert."),
        ("Citation discipline",
         "Answers point to a specific lecture deck and slide number.",
         "Trains students to verify — how good engineers actually study."),
        ("Bilingual by default",
         "EN / 中文 inline; the model matches the student's register.",
         "Cantonese students think across languages. The tool should too."),
        ("Failure-mode awareness",
         "Generates trap-analysis material, not only correct answers.",
         "Most marks are lost on three or four well-known traps."),
        ("Spaced practice loop",
         "Quiz → wrong answer → flashcard → revisit in three days.",
         "Cognitive science 101, baked into the UI rather than left to discipline."),
    ]
    base_y = Inches(3.05)
    row_h = Inches(0.78)
    cx = [Inches(0.55), Inches(3.25), Inches(8.25)]
    cw = [Inches(2.7), Inches(5.0), Inches(4.53)]
    for i, (lbl, what, why) in enumerate(rows):
        ry = base_y + i * row_h
        if i > 0:
            line(s, Inches(0.55), ry, Inches(12.78), ry,
                 color=HAIR, width_pt=0.3)
        textbox(s, cx[0], ry + Inches(0.1), cw[0], row_h, lbl,
                font=SANS_BOLD, size=11, bold=True, color=INK)
        textbox(s, cx[1], ry + Inches(0.1), cw[1], row_h, what,
                font=SANS, size=11, color=INK, line_spacing=1.3)
        textbox(s, cx[2], ry + Inches(0.1), cw[2], row_h, why,
                font=SANS, size=10, color=INK2, line_spacing=1.3)
    return s


def slide_validation(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "10 · Validation", "We did not ship on vibes.",
               "Quantitative checks panellists can re-run on the spot.")
    line(s, Inches(6.65), Inches(3.2), Inches(6.65), Inches(6.85),
         color=HAIR, width_pt=0.4)

    textbox(s, Inches(0.55), Inches(3.15), Inches(5.8), Inches(0.3),
            "INTER-AI AGREEMENT", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(0.55), Inches(3.5), Inches(5.8), Inches(1.5),
            "We posed the 25 hardest ELEC3120 questions to LearningPacer, "
            "ChatGPT-4o, and Gemini 1.5, then graded each answer rubric-by-rubric.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    textbox(s, Inches(0.55), Inches(4.6), Inches(5.8), Inches(1.5),
            "On 21 / 25 questions all three models agree on the final answer. "
            "On the remaining 4, LearningPacer's answer matches the lecture "
            "slide; the others diverge.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    textbox(s, Inches(0.55), Inches(5.9), Inches(5.8), Inches(0.3),
            "PROTOCOL ACCURACY",
            font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(0.55), Inches(6.25), Inches(5.8), Inches(0.85),
            "On 30 protocol-fact questions from past papers: LearningPacer "
            "30 / 30; ChatGPT (no RAG) 24 / 30 — all misses were silent, "
            "confidently stated wrong values.",
            font=SANS, size=11, color=INK, line_spacing=1.4)

    textbox(s, Inches(6.95), Inches(3.15), Inches(5.85), Inches(0.3),
            "WHAT THIS IS NOT", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(6.95), Inches(3.5), Inches(5.85), Inches(0.85),
            "Not a peer-reviewed user study. We are FYP students. The panel "
            "can re-run the questions live during the demo.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    textbox(s, Inches(6.95), Inches(4.5), Inches(5.85), Inches(0.3),
            "WHY IT MATTERS", font=SANS_BOLD, size=10, bold=True, color=TEAL)
    textbox(s, Inches(6.95), Inches(4.85), Inches(5.85), Inches(1.0),
            "A virtual TA's job is not to be witty. It is to be wrong less "
            "often than the alternatives — and to show its work when it is right.",
            font=SANS, size=11, color=INK, line_spacing=1.4)
    qbg = rect(s, Inches(6.95), Inches(6.0), Inches(5.85), Inches(1.05),
               fill_c=SOFT, rounded=True); qbg.line.fill.background()
    rich_textbox(s, Inches(7.15), Inches(6.1), Inches(5.55), Inches(0.95), [
        [("\u201cTwo of the three commercial chatbots failed the very first "
          "subnetting question. LearningPacer cited the slide.\u201d",
          {"font": "Cambria", "size": 12, "italic": True, "color": INK})],
        [("— Internal validation log, April 2026",
          {"size": 9, "color": MUTED})],
    ], line_spacing=1.35)
    return s


def slide_diff(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "11 · Differentiation", "Why not just use ChatGPT?",
               "An honest comparison — including where the alternatives are still better.")

    YES = ("✓", TEAL); MEH = ("∼", AMBER); NO = ("—", RED)
    headers = ["Capability", "ChatGPT (free)", "Human TA",
               "Lecture PDFs", "LearningPacer"]
    rows = [
        ("Cites the exact lecture slide",          NO,  MEH, YES, YES),
        ("Generates exam-style mock papers",        MEH, NO,  NO,  YES),
        ("Bilingual EN / 中文 reasoning",          YES, MEH, NO,  YES),
        ("Available 24/7, instantly",               YES, NO,  YES, YES),
        ("Tracks individual study progress",        NO,  NO,  NO,  YES),
        ("Refuses confidently when out of scope",   NO,  YES, NO,  YES),
        ("Renders ELEC3120-style diagrams",         NO,  MEH, YES, YES),
        ("Costs the student nothing extra",         MEH, YES, YES, YES),
    ]
    base_x = Inches(0.55); base_y = Inches(2.95)
    table_w = Inches(12.23)
    cw0 = Inches(4.23)               # cap col
    cw  = (table_w - cw0) / 4        # 4 equal data cols → 2.00"
    cx  = [base_x] + [base_x + cw0 + i * cw for i in range(4)]
    rh  = Inches(0.41)

    # header band
    band = rect(s, base_x, base_y, table_w, rh,
                fill_c=SOFT2); _no_outline(band)
    line(s, base_x, base_y + rh, base_x + table_w, base_y + rh,
         color=TEAL, width_pt=0.7)
    for i, h in enumerate(headers):
        textbox(s, cx[i] + Inches(0.1), base_y + Inches(0.1),
                (cw0 if i == 0 else cw) - Inches(0.2), rh, h,
                font=SANS_BOLD, size=11, bold=True, color=INK,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
    # rows
    for r_i, row in enumerate(rows):
        ry = base_y + (r_i + 1) * rh
        # last column highlight (LearningPacer)
        last = rect(s, cx[4], ry, cw, rh, fill_c=SOFT); _no_outline(last)
        # text
        textbox(s, cx[0] + Inches(0.1), ry + Inches(0.08),
                cw0 - Inches(0.2), rh, row[0],
                font=SANS, size=10.5, color=INK)
        for i, (sym, col) in enumerate(row[1:], start=1):
            textbox(s, cx[i], ry + Inches(0.06), cw, rh, sym,
                    font=SANS_BOLD, size=14, bold=True, color=col,
                    align=PP_ALIGN.CENTER)
        line(s, base_x, ry + rh, base_x + table_w, ry + rh,
             color=HAIR, width_pt=0.3)

    legend_y = base_y + (len(rows) + 1) * rh + Inches(0.15)
    rich_textbox(s, base_x, legend_y, table_w, Inches(0.3), [[
        ("✓ ", {"bold": True, "color": TEAL}),
        ("supported    ", {"color": INK2, "size": 10}),
        ("∼ ", {"bold": True, "color": AMBER}),
        ("partial / inconsistent    ", {"color": INK2, "size": 10}),
        ("— ", {"bold": True, "color": RED}),
        ("not supported", {"color": INK2, "size": 10}),
    ]])
    return s


def slide_risks(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "12 · Risks & mitigations",
               "We pre-empted the panel's hardest questions.")
    rows = [
        ("Hallucination",
         "Model invents a protocol detail absent from the slides.",
         "Retrieval-grounded answers; refusal mode when no source matches; trap-analysis trains students to verify."),
        ("Academic integrity",
         "Used to write graded assignments instead of to learn.",
         "No homework-completion mode; answers are explanatory and cite the lecture; institution can opt in to logging."),
        ("Cost & scaling",
         "LLM spend balloons at exam time.",
         "Two-provider routing; cached embeddings; per-user daily soft caps."),
        ("Privacy",
         "Student notes and conversation history.",
         "Supabase RLS; only the student reads their own rows; no third-party analytics on chat content."),
        ("Single-course scope",
         "Only ELEC3120 today.",
         "Knowledge base is course-scoped behind a feature flag; same architecture extends to ELEC3100, COMP3511."),
    ]
    base_x = Inches(0.55); base_y = Inches(2.75)
    cw = [Inches(2.0), Inches(3.6), Inches(6.63)]
    cx = [base_x, base_x + cw[0], base_x + cw[0] + cw[1]]
    head_h = Inches(0.40)
    row_h  = Inches(0.76)
    # footprint: 2.75 + 0.40 + 5*0.76 = 6.95"  (clears footer rule at 7.05")
    band = rect(s, base_x, base_y, Inches(12.23), head_h, fill_c=SOFT2); _no_outline(band)
    line(s, base_x, base_y + head_h, base_x + Inches(12.23), base_y + head_h,
         color=TEAL, width_pt=0.7)
    for i, h in enumerate(["Risk", "Concrete failure mode", "Mitigation in production today"]):
        textbox(s, cx[i] + Inches(0.1), base_y + Inches(0.1),
                cw[i] - Inches(0.2), head_h, h,
                font=SANS_BOLD, size=11, bold=True, color=INK)
    for r_i, (a, b, c) in enumerate(rows):
        ry = base_y + head_h + r_i * row_h
        textbox(s, cx[0] + Inches(0.1), ry + Inches(0.12),
                cw[0] - Inches(0.2), row_h, a,
                font=SANS_BOLD, size=11, bold=True, color=INK)
        textbox(s, cx[1] + Inches(0.1), ry + Inches(0.12),
                cw[1] - Inches(0.2), row_h, b,
                font=SANS, size=10.5, color=INK, line_spacing=1.3)
        textbox(s, cx[2] + Inches(0.1), ry + Inches(0.12),
                cw[2] - Inches(0.2), row_h, c,
                font=SANS, size=10.5, color=INK, line_spacing=1.3)
        line(s, base_x, ry + row_h, base_x + Inches(12.23), ry + row_h,
             color=HAIR, width_pt=0.3)
    return s


def slide_demo(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "13 · Live demo plan", "Eight minutes. Six beats.",
               "What you will see, in the order you will see it.")
    demo_timeline(s, Inches(0.55), Inches(3.1), Inches(12.23), Inches(1.6))
    rows = [
        ("00:30", "Open the live app on the conference wifi — no localhost, no mock screenshots."),
        ("02:00", "Ask a hard ELEC3120 question (TCP CUBIC vs Reno with derivation)."),
        ("03:30", "Generate a brand-new mock midterm in front of the panel; show print preview."),
        ("05:30", "Run a 5-question quiz on a topic the panel picks; watch the stats ring update."),
        ("07:00", "Open the knowledge base, glossary, and subnet calculator."),
        ("08:00", "Q&A — 25-question pack rehearsed, but unrehearsed questions welcome."),
    ]
    base_y = Inches(4.95)
    row_h = Inches(0.32)
    for i, (t, d) in enumerate(rows):
        ry = base_y + i * row_h
        textbox(s, Inches(0.55), ry, Inches(1.0), row_h, t,
                font=SANS_BOLD, size=10.5, bold=True, color=TEAL)
        textbox(s, Inches(1.65), ry, Inches(11.13), row_h, d,
                font=SANS, size=10.5, color=INK, line_spacing=1.3)
        if i < len(rows) - 1:
            line(s, Inches(0.55), ry + row_h, Inches(12.78), ry + row_h,
                 color=HAIR, width_pt=0.25)
    return s


def slide_roadmap(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    add_header(s, "14 · Roadmap", "Past the FYP, into the next semester.",
               "Ranked by what students asked for in the pilot.")
    rows = [
        ("Now (May 2026)",
         "FYP demo build · ELEC3120 only · single-tenant Supabase · one lecturer's corpus."),
        ("Q3 2026",
         "Multi-course mode · ELEC3100, COMP3511 · per-course knowledge-base flags · "
         "instructor admin panel."),
        ("Q4 2026",
         "Spaced-repetition engine that re-asks weak-topic flashcards on a real schedule; "
         "exam-week cram pack auto-generated from study history."),
        ("2027",
         "Department-licensed deployment · LTI integration with Canvas · "
         "anonymised study analytics for course coordinators."),
    ]
    base_y = Inches(3.05)
    row_h = Inches(0.85)
    for i, (label, desc) in enumerate(rows):
        ry = base_y + i * row_h
        if i > 0:
            line(s, Inches(0.55), ry, Inches(12.78), ry,
                 color=HAIR, width_pt=0.3)
        textbox(s, Inches(0.55), ry + Inches(0.12), Inches(2.6), row_h, label,
                font=SANS_BOLD, size=11.5, bold=True, color=TEAL)
        textbox(s, Inches(3.4), ry + Inches(0.12), Inches(9.4), row_h, desc,
                font=SANS, size=11, color=INK, line_spacing=1.4)
    return s


def slide_closing(page, total):
    s = prs.slides.add_slide(BLANK)
    add_chrome_normal(s, page, total)
    textbox(s, Inches(0.55), Inches(1.55), Inches(12), Inches(0.4),
            "THANK YOU",
            font=SANS_BOLD, size=12, bold=True, color=TEAL)
    textbox(s, Inches(0.55), Inches(1.95), Inches(12), Inches(2.2),
            "Questions are how\nthe tutor learns, too.",
            font=SERIF_BOLD, size=54, bold=True, color=INK,
            line_spacing=1.05)
    line(s, Inches(0.55), Inches(4.4), Inches(8.5), Inches(4.4),
         color=TEAL, width_pt=1.5)
    rows = [
        ("PROJECT",             "LearningPacer  ·  ELEC3120 Virtual TA  ·  HKUST"),
        ("LIVE DEMO",           "Available now on the conference wifi"),
        ("COMPANION ARTIFACTS", "Mock Exam · Trap Analysis · 25-Q Defence Pack"),
        ("OPEN TO",             "Code review · panel questions · unrehearsed prompts"),
    ]
    base_y = Inches(4.7); row_h = Inches(0.42)
    for i, (k, v) in enumerate(rows):
        ry = base_y + i * row_h
        textbox(s, Inches(0.55), ry, Inches(2.6), row_h, k,
                font=SANS_BOLD, size=11, bold=True, color=TEAL)
        textbox(s, Inches(3.2),  ry, Inches(9.5), row_h, v,
                font=SANS, size=13, color=INK)
    return s


# ── build ────────────────────────────────────────────────────────────────────
def build():
    BUILDERS = [
        ("cover", slide_cover),
        slide_context,
        slide_problem,
        slide_solution,
        slide_architecture,
        slide_stack,
        slide_tutor,
        slide_exam,
        slide_study,
        slide_pedagogy,
        slide_validation,
        slide_diff,
        slide_risks,
        slide_demo,
        slide_roadmap,
        slide_closing,
    ]
    total = len(BUILDERS)
    page = 1
    for entry in BUILDERS:
        if isinstance(entry, tuple) and entry[0] == "cover":
            entry[1]()
        else:
            entry(page, total)
        page += 1
    prs.save(OUT_PATH)
    return OUT_PATH, total


if __name__ == "__main__":
    path, n = build()
    kb = os.path.getsize(path) / 1024
    print(f"OK  {path}  ({kb:.1f} KB, {n} slides)")
